#!/usr/bin/env python3
"""Generate Img2CAD roadshow PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ─── Color Palette ───
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BG_SOFT    = RGBColor(0xF5, 0xF7, 0xFA)
TEXT_DARK  = RGBColor(0x0A, 0x25, 0x40)
TEXT_MID   = RGBColor(0x42, 0x54, 0x66)
TEXT_DIM   = RGBColor(0x88, 0x98, 0xAA)
BLUE       = RGBColor(0x25, 0x63, 0xEB)
PURPLE     = RGBColor(0x7C, 0x3A, 0xED)
CYAN       = RGBColor(0x08, 0x91, 0xB2)
GREEN      = RGBColor(0x16, 0xA3, 0x4A)
RED        = RGBColor(0xDC, 0x26, 0x26)
AMBER      = RGBColor(0xD9, 0x77, 0x06)
SLATE_800  = RGBColor(0x1E, 0x29, 0x3B)
BORDER     = RGBColor(0xD0, 0xD8, 0xE4)
BLUE_BG    = RGBColor(0xEF, 0xF6, 0xFF)
PURPLE_BG  = RGBColor(0xFA, 0xF5, 0xFF)
CYAN_BG    = RGBColor(0xEC, 0xFE, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ─── Helper Functions ───

def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, radius=None, shadow=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    if radius and hasattr(shape, 'adjustments'):
        r = int(radius / Inches(1) * 100) if isinstance(radius, int) else 5000
    if shadow:
        shape.shadow.inherit = False
    return shape

def add_text_box(slide, left, top, width, height, text="", font_size=14, color=TEXT_DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei', line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(2)
    if line_spacing != 1.0:
        p.line_spacing = Pt(int(font_size * line_spacing))
    return txBox, tf

def add_rich_text_box(slide, left, top, width, height):
    """Returns (text_frame, shape) for manual paragraph building."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf, txBox

def add_paragraph(tf, text, font_size=14, color=TEXT_DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei', space_after=4, first=False):
    if first:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(space_after)
    return p

def add_section_num(slide, num_text):
    add_text_box(slide, Inches(10.8), Inches(6.2), Inches(2.2), Inches(1.0),
                 num_text, font_size=140, color=RGBColor(0xE0, 0xE8, 0xF8), bold=True,
                 alignment=PP_ALIGN.RIGHT)

def add_kicker(slide, text, left=Inches(0.8), top=Inches(0.5)):
    add_text_box(slide, left, top, Inches(5), Inches(0.4), text, font_size=13, color=BLUE, bold=True)

def add_title(slide, text, left=Inches(0.8), top=Inches(1.0), width=Inches(10)):
    add_text_box(slide, left, top, width, Inches(1.0), text, font_size=38, color=TEXT_DARK, bold=True)

def add_subtitle(slide, text, left=Inches(0.8), top=Inches(1.7)):
    add_text_box(slide, left, top, Inches(10), Inches(0.5), text, font_size=14, color=TEXT_DIM)

def add_card(slide, left, top, width, height, title, body, icon=None, accent=False):
    fill = BLUE_BG if accent else WHITE
    border = RGBColor(0xBF, 0xD0, 0xF0) if accent else BORDER
    shape = add_rect(slide, left, top, width, height, fill_color=fill, border_color=border)
    y = top + Inches(0.2)
    if icon:
        add_text_box(slide, left + Inches(0.25), y, Inches(0.5), Inches(0.5), icon, font_size=22)
        y += Inches(0.4)
    else:
        y += Inches(0.1)
    add_text_box(slide, left + Inches(0.25), y, width - Inches(0.5), Inches(0.35), title, font_size=16, color=TEXT_DARK, bold=True)
    y += Inches(0.35)
    add_text_box(slide, left + Inches(0.25), y, width - Inches(0.5), height - Inches(0.7) - (y - top), body, font_size=12, color=TEXT_DIM)

def add_metric_card(slide, left, top, width, height, number, label, desc=""):
    shape = add_rect(slide, left, top, width, height, fill_color=WHITE, border_color=BORDER)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.25), width - Inches(0.4), Inches(0.7),
                 number, font_size=50, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), top + Inches(1.0), width - Inches(0.4), Inches(0.3),
                 label, font_size=14, color=TEXT_MID, alignment=PP_ALIGN.CENTER)
    if desc:
        add_text_box(slide, left + Inches(0.2), top + Inches(1.35), width - Inches(0.4), Inches(0.5),
                     desc, font_size=11, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

def add_pill(slide, left, top, text, accent=False):
    w = Inches(len(text) * 0.12 + 0.4)
    fill = RGBColor(0xDB, 0xE4, 0xFA) if accent else BG_SOFT
    fcolor = BLUE if accent else TEXT_MID
    shape = add_rect(slide, left, top, w, Inches(0.32), fill_color=fill)
    add_text_box(slide, left, top, w, Inches(0.32), text, font_size=11, color=fcolor, bold=True, alignment=PP_ALIGN.CENTER)

def add_code_block(slide, left, top, width, height, lines):
    """Add a dark code block with syntax-ish coloring."""
    shape = add_rect(slide, left, top, width, height, fill_color=SLATE_800)
    tf, _ = add_rich_text_box(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), height - Inches(0.3))
    for i, (text, color) in enumerate(lines):
        add_paragraph(tf, text, font_size=10, color=color, font_name='Consolas', space_after=1, first=(i==0))

def add_flow_arrow(slide, left, top):
    add_text_box(slide, left, top, Inches(0.35), Inches(0.35), "▸", font_size=18, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)

def add_flow_step(slide, left, top, icon, label, sub, highlight=False):
    border = BLUE if highlight else BORDER
    shape = add_rect(slide, left, top, Inches(1.45), Inches(1.3), fill_color=BG_SOFT, border_color=border)
    add_text_box(slide, left + Inches(0.05), top + Inches(0.1), Inches(1.35), Inches(0.4), icon, font_size=22, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.05), top + Inches(0.5), Inches(1.35), Inches(0.3), label, font_size=11, color=TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.05), top + Inches(0.8), Inches(1.35), Inches(0.4), sub, font_size=9, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

def add_footer(slide, left_text, slide_num, total=9):
    add_text_box(slide, Inches(0.8), Inches(7.1), Inches(5), Inches(0.3), left_text, font_size=10, color=TEXT_DIM)
    add_text_box(slide, Inches(11.8), Inches(7.1), Inches(1), Inches(0.3),
                 f"{slide_num} / {total}", font_size=10, color=TEXT_DIM, alignment=PP_ALIGN.RIGHT)

def add_cover_blob(slide):
    # Gradient blobs approximated with shapes
    blob1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-1.0), Inches(4.5), Inches(4.5))
    blob1.fill.solid(); blob1.fill.fore_color.rgb = BLUE_BG
    blob1.line.fill.background()
    blob1.rotation = 15
    blob2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.0), Inches(5.5), Inches(3.2), Inches(3.2))
    blob2.fill.solid(); blob2.fill.fore_color.rgb = PURPLE_BG
    blob2.line.fill.background()

# ═══════════════════════════════════════════
# SLIDE 1: COVER
# ═══════════════════════════════════════════
s = add_blank_slide()
add_bg(s, WHITE)
# gradient bg approximation
bg_shape = add_rect(s, 0, 0, W, H, fill_color=RGBColor(0xF0, 0xF6, 0xFE))
add_cover_blob(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(4), Inches(0.4), "● Img2CAD", font_size=16, color=TEXT_DARK, bold=True)
add_text_box(s, Inches(0.8), Inches(2.2), Inches(5), Inches(0.4), "北京邮电大学 2025 雏雁计划 · 路演展示", font_size=14, color=BLUE, bold=True)
add_text_box(s, Inches(0.8), Inches(2.8), Inches(11), Inches(1.5),
             "一拍即得\n手机照片 → 三维 CAD 模型", font_size=58, color=TEXT_DARK, bold=True)
add_text_box(s, Inches(0.8), Inches(4.5), Inches(10), Inches(0.8),
             "基于 VLM 辅助条件分解的单视图图像三维 CAD 逆向工程\n技术基础：SIGGRAPH Asia 2025 (CCF-A) · Stanford / NVIDIA / Peking Univ. · MIT 开源",
             font_size=16, color=TEXT_MID)
# Pills
add_pill(s, Inches(0.8), Inches(5.6), "📸 手机拍照输入", accent=True)
add_pill(s, Inches(3.0), Inches(5.6), "⚡ 60-120s 出结果", accent=True)
add_pill(s, Inches(5.2), Inches(5.6), "📦 OBJ + STEP 输出", accent=True)
add_footer(s, "Img2CAD Team", 1)

# ═══════════════════════════════════════════
# SLIDE 2: PAIN POINT
# ═══════════════════════════════════════════
s = add_blank_slide()
add_bg(s, WHITE)
add_section_num(s, "01")
add_kicker(s, "行业痛点")
add_title(s, "CAD 建模：高门槛、低效率")

cards_data = [
    ("⏳", "耗时长", "复杂模型需数小时至数天\n人工操作 SolidWorks / Fusion 360"),
    ("💰", "成本高", "单模型 50-200 美元\n中小企业难以负担专业团队"),
    ("🔧", "门槛高", "需数年专业培训\n将普通用户拒之门外"),
]
for i, (icon, title, body) in enumerate(cards_data):
    x = Inches(0.8 + i * 4.1)
    add_card(s, x, Inches(2.3), Inches(3.8), Inches(2.0), title, body, icon=icon)

add_card(s, Inches(0.8), Inches(4.6), Inches(11.7), Inches(1.5),
         "🎯 AI 时代的核心挑战",
         "CAD 模型不同于普通 3D 网格 — 必须同时满足三大约束：精确几何参数 + 结构化构造历史 + 拓扑正确性。\n现有端到端深度学习方法均无法同时兼顾三者。",
         accent=True)
add_footer(s, "Img2CAD · 痛点分析", 2)

# ═══════════════════════════════════════════
# SLIDE 3: SOLUTION
# ═══════════════════════════════════════════
s = add_blank_slide()
add_bg(s, WHITE)
add_section_num(s, "02")
add_kicker(s, "解决方案")
add_title(s, "条件分解：让 AI 各司其职")
add_subtitle(s, "核心思想：将复杂的 Image→CAD 跨模态映射，分解为两个条件独立的子任务")

stage1_text = (
    "Llama-3.2-11B-Vision + QLoRA 微调\n\n"
    "✓ 从照片理解语义：这是椅子还是桌子？由哪些部件组成？\n"
    "✓ 预测 CAD 命令序列：部件名 · 草图曲线类型 · 布尔操作 · 挤出方式\n"
    "✓ 发挥 VLM 的天然优势 — 视觉语义理解与推理"
)
stage2_text = (
    "TrAssembler · Gaussian Mixture Flow 扩散模型\n\n"
    "✓ 以 Stage 1 结构 + 输入图像为条件，预测精确的连续参数值\n"
    "✓ 32 混合分量建模多模态参数分布 · FlowEulerODE 采样器\n"
    "✓ 发挥扩散模型的天然优势 — 高精度数值回归"
)
add_card(s, Inches(0.8), Inches(2.4), Inches(5.7), Inches(3.2), "🧠  Stage 1 · 离散结构预测", stage1_text, accent=True)
add_card(s, Inches(6.8), Inches(2.4), Inches(5.7), Inches(3.2), "📐  Stage 2 · 连续参数回归", stage2_text, accent=True)

add_card(s, Inches(0.8), Inches(5.9), Inches(11.7), Inches(0.8),
         "数学基础",
         "P(CAD | Image) = P(Structure | Image) × P(Parameters | Structure, Image)  — 贝叶斯条件分解，数学上严格成立\n论文发表于 SIGGRAPH Asia 2025（CCF-A 类，计算机图形学顶级会议）",
         accent=True)
add_footer(s, "Img2CAD · 解决方案", 3)

# ═══════════════════════════════════════════
# SLIDE 4: ARCHITECTURE
# ═══════════════════════════════════════════
s = add_blank_slide()
add_bg(s, WHITE)
add_section_num(s, "03")
add_kicker(s, "技术架构")
add_title(s, "端到端推理流水线")
add_subtitle(s, "核心实现：infer_single.py（~900 行 Python，13 个解析函数）— 全自动串联所有阶段")

flow_steps = [
    ("📸", "输入图片", "手机实拍\nJPG / PNG", False),
    ("✂️", "预处理", "rembg ISNet\n形态学修复", False),
    ("🧠", "Stage 1", "Llama-3.2-11B\n+ LoRA Adapter", True),
    ("📐", "Stage 2", "TrAssembler\nGMFlow 扩散", False),
    ("🔧", "实体生成", "pythonOCC\n几何引擎", False),
    ("📦", "导出", "OBJ + STEP\n+ PNG 预览", False),
]
fx = Inches(0.5)
for i, (icon, label, sub, hl) in enumerate(flow_steps):
    x = fx + i * Inches(2.05)
    add_flow_step(s, x, Inches(2.4), icon, label, sub, highlight=hl)
    if i < len(flow_steps) - 1:
        add_flow_arrow(s, x + Inches(1.48), Inches(2.85))

# Code block
code_lines = [
    ("# infer_single.py — 端到端单图推理入口", RGBColor(0x94, 0xA3, 0xB8)),
    ("def infer_single(image_path, category=\"chair\"):", RGBColor(0x7D, 0xD3, 0xFC)),
    ("    img = preprocess_phone_photo(image_path)", RGBColor(0xF0, 0xAB, 0xFC)),
    ("    # Stage 1: VLM → 离散 CAD 结构", RGBColor(0x94, 0xA3, 0xB8)),
    ("    model = load_llama_4bit(adapter_path)", RGBColor(0xF0, 0xAB, 0xFC)),
    ("    text = llama_generate(img)          # ~80 行", RGBColor(0xF0, 0xAB, 0xFC)),
    ("    h5   = text_to_h5(text)            # ~270 行，13 个解析函数", RGBColor(0xF0, 0xAB, 0xFC)),
    ("    del model; torch.cuda.empty_cache() # 释放显存", RGBColor(0x94, 0xA3, 0xB8)),
    ("    # Stage 2: GMFlow → 连续参数", RGBColor(0x94, 0xA3, 0xB8)),
    ("    model = load_trassembler(ckpt_path)", RGBColor(0xF0, 0xAB, 0xFC)),
    ("    params = diffusion_sample(img, h5) # ~120 行", RGBColor(0xF0, 0xAB, 0xFC)),
    ("    return vec2CADsolid(params)      # → .obj + .step", RGBColor(0x94, 0xA3, 0xB8)),
]
add_code_block(s, Inches(0.8), Inches(4.1), Inches(5.7), Inches(2.8), code_lines)

# GPU sidebar
add_card(s, Inches(6.8), Inches(4.1), Inches(5.7), Inches(2.8),
         "⚡ GPU 显存调度策略",
         "串行加载 → 即时释放\nStage 1 推理完成后立即 del VLM 并 torch.cuda.empty_cache()\n\n"
         "❌ 同时加载两模型：44GB+\n"
         "✓ 串行调度 + 4-bit QLoRA：仅需 22GB\n\n"
         "消费级 RTX 3090 / 4090 (24GB) 即可运行 —\nStage 1 ~22GB，Stage 2 ~4GB",
         accent=True)
add_footer(s, "Img2CAD · 技术架构", 4)

# ═══════════════════════════════════════════
# SLIDE 5: ORIGINAL CONTRIBUTIONS
# ═══════════════════════════════════════════
s = add_blank_slide()
add_bg(s, WHITE)
add_section_num(s, "04")
add_kicker(s, "原创贡献")
add_title(s, "从"论文代码"到"可用产品"")
add_subtitle(s, "论文提供理论基础与模块化训练/评估脚本 — 本项目完成 7 项关键工程化原创贡献")

contribs = [
    ("1", "端到端推理流水线 · infer_single.py ~900 行", "13 个处理函数全自动串联：预处理 → Stage1 → 文本解析 → Stage2 → 实体生成 → 多格式导出"),
    ("2", "手机照片域适配预处理 · ~80 行", "rembg(ISNet) 背景移除 → 形态学闭运算修复薄结构 → 最大连通分量去噪 → 白底居中"),
    ("3", "VLM 输出鲁棒解析器 · ~270 行", "5 层容错：filter_lines → divide_blocks → filter_blocks → 变量求值引擎 → arc 顺序纠正 + CLIP 语义纠错。解析成功率 ~60% → 90%+"),
    ("4", "Web 服务化架构 · server/ ~660 行", "FastAPI + WebSocket · UUID 作业状态机(queued→processing→done/error) · 按类别模型缓存(热启动~5s) · 自动过期清理"),
    ("5", "React 前端 · frontend/src/ 6 组件", "UploadZone · CategorySelector · ProgressPanel(WebSocket 实时进度) · ResultPanel(预览+下载) · ErrorPanel · Header"),
    ("6", "公网一键部署 · start.sh / stop.sh", "cpolar VIP 固定子域名 img2cad.cpolar.cn · 一键启动(清理→隧道→服务→健康检查) · 一键停止"),
    ("7", "GPU 显存极致优化", "4-bit NF4 双重量化 (QLoRA) + Stage 1→2 即时释放：总需求 44GB → 22GB，消费级显卡可运行全流程"),
]
for i, (num, title, desc) in enumerate(contribs):
    col = 0 if i % 2 == 0 else 1
    row = i // 2
    x = Inches(0.8 + col * 6.1)
    y = Inches(2.5 + row * 1.05)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y + Inches(0.05), Inches(0.35), Inches(0.35))
    circle.fill.solid(); circle.fill.fore_color.rgb = BLUE
    circle.line.fill.background()
    tf = circle.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = num; p.font.size = Pt(13); p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    # Text
    add_text_box(s, x + Inches(0.5), y, Inches(5.3), Inches(0.3), title, font_size=13, color=TEXT_DARK, bold=True)
    add_text_box(s, x + Inches(0.5), y + Inches(0.3), Inches(5.3), Inches(0.6), desc, font_size=10, color=TEXT_DIM)

add_footer(s, "Img2CAD · 原创贡献", 5)

# ═══════════════════════════════════════════
# SLIDE 6: METRICS
# ═══════════════════════════════════════════
s = add_blank_slide()
add_bg(s, WHITE)
add_section_num(s, "05")
add_kicker(s, "关键指标")
add_title(s, "数据说话")

add_metric_card(s, Inches(0.8), Inches(2.5), Inches(2.7), Inches(2.0), "3", "家具类别", "椅子 · 桌子 · 柜子")
add_metric_card(s, Inches(3.8), Inches(2.5), Inches(2.7), Inches(2.0), "90%+", "解析成功率", "VLM 输出一次解析")
add_metric_card(s, Inches(6.8), Inches(2.5), Inches(2.7), Inches(2.0), "22 GB", "GPU 显存", "消费级显卡可运行")
add_metric_card(s, Inches(9.8), Inches(2.5), Inches(2.7), Inches(2.0), "60-120s", "推理时间", "上传到下载全流程")

add_card(s, Inches(0.8), Inches(4.8), Inches(5.7), Inches(1.7),
         "📐 输出格式",
         "OBJ 网格格式  ·  STEP 工业标准  ·  PNG 渲染预览\n\nSTEP 可直接导入 SolidWorks / Fusion 360 / AutoCAD 进行二次编辑")
add_card(s, Inches(6.8), Inches(4.8), Inches(5.7), Inches(1.7),
         "⚡ 效率对比",
         "手工建模：数小时 → Img2CAD：60-120s\n效率提升数十倍\n成本从 50-200 USD/模型 → 0.1 USD/次")
add_footer(s, "Img2CAD · 关键指标", 6)

# ═══════════════════════════════════════════
# SLIDE 7: VALUE & IMPACT
# ═══════════════════════════════════════════
s = add_blank_slide()
add_bg(s, WHITE)
add_section_num(s, "06")
add_kicker(s, "价值与影响")
add_title(s, "技术落地的三重价值")

values = [
    ("🌍", "社会价值", "降低三维建模门槛 — 手机拍照即可获得 CAD 模型\nSTEM 教育普惠 — 三维设计教学\n文化遗产三维数字化保护\n电商 3D 商品展示 — 降低退货率和物流浪费"),
    ("💰", "经济效益", "减少 60-80% 基础建模工时\n10 人团队年节约 50-100 万人力成本\n云端推理仅 0.1-0.2 美元/次\n对比人工：50-200 美元 → 0.1 美元"),
    ("🏭", "战略意义", '契合 "中国制造 2025" 战略方向\n呼应 "工业 4.0" 智能制造趋势\nAI 赋能传统制造业的典型范式\n连接学术前沿与产业需求的关键桥梁'),
]
for i, (icon, title, body) in enumerate(values):
    x = Inches(0.8 + i * 4.1)
    add_card(s, x, Inches(2.5), Inches(3.8), Inches(3.0), title, body, icon=icon, accent=True)

add_card(s, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.9),
         "核心价值",
         "将 SIGGRAPH Asia 2025 前沿研究转化为可实际部署的工程系统 —"
         "从\"论文代码\"到\"可用产品\"的工程化过程本身，即是本项目最核心的社会价值",
         accent=True)
add_footer(s, "Img2CAD · 价值与影响", 7)

# ═══════════════════════════════════════════
# SLIDE 8: ROADMAP
# ═══════════════════════════════════════════
s = add_blank_slide()
add_bg(s, WHITE)
add_section_num(s, "07")
add_kicker(s, "发展路线")
add_title(s, "从家具到万物建模")

roadmap = [
    ("✅ 已完成", "当前阶段", "3 类家具端到端推理\nWeb 服务 + React 前端\nimg2cad.cpolar.cn\nMIT 开源 · HuggingFace", True, GREEN),
    ("🔜 短期", "品类扩展", "机械零件 · 齿轮轴承法兰\n建筑构件 · 门窗楼梯柱\n更多家具品类\n训练数据扩展管线", False, BLUE),
    ("🎯 中长期", "能力跃升", "多视角融合推理提升精度\n视频流实时 CAD 重建\n移动端轻量化部署\n\"拍照即建模\"通用平台", False, PURPLE),
]
for i, (badge, title, body, accent, badge_color) in enumerate(roadmap):
    x = Inches(0.8 + i * 4.1)
    fill = BLUE_BG if accent else WHITE
    border = RGBColor(0xBF, 0xD0, 0xF0) if accent else BORDER
    shape = add_rect(s, x, Inches(2.5), Inches(3.8), Inches(3.0), fill_color=fill, border_color=border)
    # Badge
    add_text_box(s, x + Inches(0.3), Inches(2.75), Inches(3.2), Inches(0.3), badge, font_size=12, color=badge_color, bold=True)
    add_text_box(s, x + Inches(0.3), Inches(3.15), Inches(3.2), Inches(0.35), title, font_size=20, color=TEXT_DARK, bold=True)
    add_text_box(s, x + Inches(0.3), Inches(3.6), Inches(3.2), Inches(1.6), body, font_size=12, color=TEXT_DIM)

add_card(s, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.9),
         "长期愿景",
         "打造 \"拍照即建模\" 的通用三维逆向工程平台，服务数字孪生 · 智能制造 · AR/VR · 机器人仿真等万亿级市场场景",
         accent=True)
add_footer(s, "Img2CAD · 发展路线", 8)

# ═══════════════════════════════════════════
# SLIDE 9: THANK YOU
# ═══════════════════════════════════════════
s = add_blank_slide()
add_bg(s, WHITE)
bg2 = add_rect(s, 0, 0, W, H, fill_color=RGBColor(0xF0, 0xF6, 0xFE))
add_cover_blob(s)

add_text_box(s, Inches(0.8), Inches(2.8), Inches(11.5), Inches(1.0),
             "感谢聆听", font_size=64, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(0.8), Inches(3.8), Inches(11.5), Inches(0.5),
             "欢迎提问与交流", font_size=20, color=TEXT_MID, alignment=PP_ALIGN.CENTER)

# CTA box
cta_shape = add_rect(s, Inches(3.0), Inches(4.6), Inches(7.3), Inches(1.8), fill_color=BLUE)
add_text_box(s, Inches(3.3), Inches(4.75), Inches(6.7), Inches(0.4),
             "🔗 在线体验 & 开源地址", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(3.3), Inches(5.2), Inches(6.7), Inches(0.4),
             "http://img2cad.cpolar.cn", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(3.3), Inches(5.6), Inches(6.7), Inches(0.35),
             "公网可直接访问，手机拍照即刻体验", font_size=12, color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(3.3), Inches(6.0), Inches(6.7), Inches(0.35),
             "GitHub (MIT)  ·  HuggingFace 模型  ·  arXiv: 2408.01437", font_size=12, color=RGBColor(0xCC, 0xD8, 0xF8), alignment=PP_ALIGN.CENTER)

add_text_box(s, Inches(0.8), Inches(6.7), Inches(11.5), Inches(0.4),
             'You et al. "Img2CAD: Reverse Engineering 3D CAD Models from Images through VLM-Assisted Conditional Factorization"',
             font_size=9, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(0.8), Inches(6.9), Inches(11.5), Inches(0.3),
             "SIGGRAPH Asia 2025 (CCF-A) · Stanford / NVIDIA / Peking Univ. / DEVCOM ARL",
             font_size=9, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
add_footer(s, "Img2CAD · 北京邮电大学雏雁计划", 9)

# ─── Save ───
output_path = "/root/autodl-tmp/Img2CAD/Img2CAD_路演PPT.pptx"
prs.save(output_path)
print(f"PPTX saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
